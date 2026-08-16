import json
import os
import uuid
from datetime import datetime
from io import BytesIO

import streamlit as st
import streamlit.components.v1 as components
from dotenv import load_dotenv
from openai import OpenAI
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


# ==================================================
# OPSÆTNING
# ==================================================

load_dotenv()
client = OpenAI()

FAKTAFIL = "facts.json"
CHATFIL = "chat_history.json"
OPGAVEFIL = "tasks.json"
PROFILFIL = "partner_profile.txt"


# ==================================================
# GENEREL FILHÅNDTERING
# ==================================================

def hent_json(filnavn, standard):
    if os.path.exists(filnavn):
        try:
            with open(filnavn, "r", encoding="utf-8") as fil:
                return json.load(fil)
        except (json.JSONDecodeError, OSError):
            return standard

    return standard


def gem_json(filnavn, data):
    with open(filnavn, "w", encoding="utf-8") as fil:
        json.dump(
            data,
            fil,
            ensure_ascii=False,
            indent=2,
        )


def hent_partnerprofil():
    if os.path.exists(PROFILFIL):
        with open(PROFILFIL, "r", encoding="utf-8") as fil:
            return fil.read()

    return (
        "Du er AK-AI-Partner. "
        "Du hjælper brugeren klart, praktisk og på dansk."
    )


partnerprofil = hent_partnerprofil()


# ==================================================
# OPGAVER
# ==================================================

def ny_opgave(
    titel,
    beskrivelse="",
    ansvarlig="Ikke angivet",
    deadline="Ikke angivet",
    prioritet="MELLEM",
    kilde="Manuelt oprettet",
):
    return {
        "id": str(uuid.uuid4()),
        "titel": titel.strip(),
        "beskrivelse": beskrivelse.strip(),
        "ansvarlig": ansvarlig.strip() or "Ikke angivet",
        "deadline": deadline.strip() or "Ikke angivet",
        "prioritet": prioritet,
        "status": "ÅBEN",
        "kilde": kilde,
        "oprettet": datetime.now().isoformat(
            timespec="seconds"
        ),
        "afsluttet": None,
    }


def gem_opgaver():
    gem_json(
        OPGAVEFIL,
        st.session_state.tasks,
    )


def opgave_findes(titel, ansvarlig, deadline):
    titel = titel.strip().lower()
    ansvarlig = ansvarlig.strip().lower()
    deadline = deadline.strip().lower()

    for opgave in st.session_state.tasks:
        if opgave["status"] != "ÅBEN":
            continue

        if (
            opgave["titel"].strip().lower() == titel
            and opgave["ansvarlig"].strip().lower() == ansvarlig
            and opgave["deadline"].strip().lower() == deadline
        ):
            return True

    return False


def opret_ai_opgave(data, kilde):
    if opgave_findes(
        data["titel"],
        data["ansvarlig"],
        data["deadline"],
    ):
        return None

    opgave = ny_opgave(
        titel=data["titel"],
        beskrivelse=data["beskrivelse"],
        ansvarlig=data["ansvarlig"],
        deadline=data["deadline"],
        prioritet=data["prioritet"],
        kilde=kilde,
    )

    st.session_state.tasks.append(opgave)
    gem_opgaver()

    return opgave


def afslut_opgave(opgave_id):
    for opgave in st.session_state.tasks:
        if opgave["id"] == opgave_id:
            opgave["status"] = "AFSLUTTET"
            opgave["afsluttet"] = datetime.now().isoformat(
                timespec="seconds"
            )
            break

    gem_opgaver()


def genaabn_opgave(opgave_id):
    for opgave in st.session_state.tasks:
        if opgave["id"] == opgave_id:
            opgave["status"] = "ÅBEN"
            opgave["afsluttet"] = None
            break

    gem_opgaver()


def slet_opgave(opgave_id):
    st.session_state.tasks = [
        opgave
        for opgave in st.session_state.tasks
        if opgave["id"] != opgave_id
    ]

    gem_opgaver()


def prioritet_sortering(opgave):
    orden = {
        "HØJ": 0,
        "MELLEM": 1,
        "LAV": 2,
    }

    return orden.get(
        opgave.get("prioritet", "MELLEM"),
        1,
    )


# ==================================================
# AI: OPGAVEDETEKTION
# ==================================================

def analyser_for_opgave(besked):
    dagens_dato = datetime.now().strftime("%Y-%m-%d")

    response = client.responses.create(
        model="gpt-5.4-mini",
        instructions=(
            "Du vurderer, om brugerens besked indeholder "
            "en konkret arbejdsopgave.\n\n"

            "Regler:\n"
            "CLEAR = nogen skal tydeligt gøre noget.\n"
            "UNCERTAIN = det kunne være en opgave, "
            "men intentionen er uklar.\n"
            "NONE = ingen opgave.\n\n"

            "Opfind aldrig ansvarlig eller deadline.\n"
            "Hvis ansvarlig ikke fremgår, brug 'Ikke angivet'.\n"
            "Hvis deadline ikke fremgår, brug 'Ikke angivet'.\n\n"

            f"Dagens dato er {dagens_dato}.\n"
            "Tydelige relative datoer som 'i morgen' "
            "eller en entydig ugedag må omsættes til dato.\n\n"

            "Prioritet skal være HØJ, MELLEM eller LAV. "
            "MELLEM er standard.\n\n"

            "Svar KUN med gyldig JSON:\n"
            "{\n"
            '  "status": "CLEAR",\n'
            '  "titel": "kort opgavetitel",\n'
            '  "beskrivelse": "kort beskrivelse",\n'
            '  "ansvarlig": "navn eller Ikke angivet",\n'
            '  "deadline": "dato/tekst eller Ikke angivet",\n'
            '  "prioritet": "MELLEM"\n'
            "}"
        ),
        input=besked,
    )

    tekst = response.output_text.strip()

    try:
        data = json.loads(tekst)
    except json.JSONDecodeError:
        return {
            "status": "NONE",
            "titel": "",
            "beskrivelse": "",
            "ansvarlig": "Ikke angivet",
            "deadline": "Ikke angivet",
            "prioritet": "MELLEM",
        }

    status = str(
        data.get("status", "NONE")
    ).upper()

    if status not in [
        "CLEAR",
        "UNCERTAIN",
        "NONE",
    ]:
        status = "NONE"

    prioritet = str(
        data.get("prioritet", "MELLEM")
    ).upper()

    if prioritet not in [
        "HØJ",
        "MELLEM",
        "LAV",
    ]:
        prioritet = "MELLEM"

    return {
        "status": status,
        "titel": str(
            data.get("titel", "")
        ).strip(),
        "beskrivelse": str(
            data.get("beskrivelse", "")
        ).strip(),
        "ansvarlig": str(
            data.get(
                "ansvarlig",
                "Ikke angivet",
            )
        ).strip() or "Ikke angivet",
        "deadline": str(
            data.get(
                "deadline",
                "Ikke angivet",
            )
        ).strip() or "Ikke angivet",
        "prioritet": prioritet,
    }


# ==================================================
# HUKOMMELSE
# ==================================================

def foreslaa_hukommelse(besked):
    response = client.responses.create(
        model="gpt-5.4-mini",
        instructions=(
            "Vurder om beskeden indeholder en stabil "
            "oplysning, der er værd at huske. "
            "Eksempler er præferencer, mål, projekter "
            "og arbejdsmetoder. "
            "Gem aldrig passwords, API-nøgler eller hemmeligheder. "
            "Hvis intet bør huskes, svar præcis INGEN. "
            "Ellers svar kun med én kort sætning."
        ),
        input=besked,
    )

    forslag = response.output_text.strip()

    if forslag.upper() == "INGEN":
        return None

    return forslag


# ==================================================
# DOKUMENTLÆSNING
# ==================================================

def laes_pdf(fil):
    reader = PdfReader(
        BytesIO(fil.getvalue())
    )

    dele = []

    for nummer, side in enumerate(
        reader.pages,
        start=1,
    ):
        tekst = side.extract_text()

        if tekst:
            dele.append(
                f"--- PDF side {nummer} ---\n"
                f"{tekst}"
            )

    return "\n\n".join(dele)


def laes_word(fil):
    document = Document(
        BytesIO(fil.getvalue())
    )

    dele = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            dele.append(
                paragraph.text
            )

    for tabel in document.tables:
        for row in tabel.rows:
            dele.append(
                " | ".join(
                    celle.text.strip()
                    for celle in row.cells
                )
            )

    return "\n".join(dele)


def laes_excel(fil):
    workbook = load_workbook(
        BytesIO(fil.getvalue()),
        data_only=True,
        read_only=True,
    )

    dele = []

    for sheet in workbook.worksheets:
        dele.append(
            f"--- Excel-ark: "
            f"{sheet.title} ---"
        )

        for row in sheet.iter_rows(
            values_only=True
        ):
            vaerdier = [
                ""
                if celle is None
                else str(celle)
                for celle in row
            ]

            if any(vaerdier):
                dele.append(
                    " | ".join(vaerdier)
                )

    return "\n".join(dele)


def laes_powerpoint(fil):
    presentation = Presentation(
        BytesIO(fil.getvalue())
    )

    dele = []

    for nummer, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        dele.append(
            f"--- PowerPoint slide "
            f"{nummer} ---"
        )

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                tekst = shape.text.strip()

                if tekst:
                    dele.append(tekst)

    return "\n".join(dele)


def laes_tekstfil(fil):
    data = fil.getvalue()

    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode(
            "latin-1",
            errors="replace",
        )


def laes_vedhaeftning(fil):
    navn = fil.name.lower()

    if navn.endswith(".pdf"):
        return laes_pdf(fil)

    if navn.endswith(".docx"):
        return laes_word(fil)

    if navn.endswith(".xlsx"):
        return laes_excel(fil)

    if navn.endswith(".pptx"):
        return laes_powerpoint(fil)

    if navn.endswith(".txt"):
        return laes_tekstfil(fil)

    raise ValueError(
        f"Filtypen understøttes ikke: "
        f"{fil.name}"
    )


# ==================================================
# MAIL INTELLIGENCE
# ==================================================

def analyser_mail(mailtekst, dokumenter):
    dokumenttekst = []

    for dokument in dokumenter:
        dokumenttekst.append(
            f"===== VEDHÆFTNING: "
            f"{dokument['navn']} =====\n"
            f"{dokument['tekst']}"
        )

    samlet_input = (
        f"EMAIL:\n{mailtekst}\n\n"
        "VEDHÆFTNINGER:\n"
        + "\n\n".join(
            dokumenttekst
        )
    )

    response = client.responses.create(
        model="gpt-5.4-mini",
        instructions=(
            f"{partnerprofil}\n\n"

            "Analyser mailen og alle vedhæftninger "
            "som ét samlet materiale. "
            "Opfind aldrig oplysninger.\n\n"

            "Brug disse overskrifter:\n"
            "## Kort resume\n"
            "## Vigtige punkter\n"
            "## Handlinger\n"
            "## Deadlines\n"
            "## Afventer\n"
            "## Vedhæftninger\n"
            "## Risici og opmærksomhedspunkter\n"
            "## Prioritet\n\n"

            "Prioritet skal være "
            "HØJ, MELLEM eller LAV."
        ),
        input=samlet_input,
    )

    return response.output_text


# ==================================================
# STREAMLIT SETUP
# ==================================================

st.set_page_config(
    page_title="AK-AI-Partner",
    page_icon="🌿",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ==================================================
# DESIGN
# ==================================================

st.markdown(
    """
    <style>

    .stApp {
        background:
            radial-gradient(
                circle at 10% 0%,
                rgba(20,255,95,0.10),
                transparent 28%
            ),
            radial-gradient(
                circle at 90% 0%,
                rgba(255,205,25,0.07),
                transparent 25%
            ),
            linear-gradient(
                180deg,
                #07100b 0%,
                #0c1610 52%,
                #070a08 100%
            );

        color: #f6faf7;
    }


    .block-container {
        max-width: 1450px;
        padding-top: 2rem;
        padding-bottom: 4rem;
    }


    section[data-testid="stSidebar"] {
        background:
            linear-gradient(
                180deg,
                #050a07 0%,
                #0a1c11 60%,
                #070a08 100%
            );

        border-right:
            1px solid rgba(50,255,110,0.13);
    }


    h1 {
        color: #f5fff7;
        font-weight: 850;
        letter-spacing: -0.045em;
    }


    h2,
    h3 {
        color: #f1faf3;
        font-weight: 750;
        letter-spacing: -0.025em;
    }


    .ak-brand {
        font-size: 0.82rem;
        font-weight: 800;
        letter-spacing: 0.20em;
        color: #39ef73;
        margin-bottom: 0.4rem;
    }


    .ak-hero {
        font-size: 2.5rem;
        line-height: 1.05;
        font-weight: 850;
        letter-spacing: -0.045em;
        color: #f5fff7;
        margin-bottom: 0.45rem;
    }


    .ak-subtitle {
        font-size: 1rem;
        color: #aebdb3;
        margin-bottom: 1.35rem;
    }


    .reggae-line {
        height: 5px;
        border-radius: 999px;
        margin-bottom: 1.8rem;

        background:
            linear-gradient(
                90deg,
                #1dcc5b 0%,
                #1dcc5b 33%,
                #f1c520 33%,
                #f1c520 66%,
                #d64138 66%,
                #d64138 100%
            );

        box-shadow:
            0 0 25px
            rgba(39,255,102,0.16);
    }


    .cockpit-card {
        background:
            linear-gradient(
                145deg,
                rgba(14,31,20,0.94),
                rgba(8,18,12,0.94)
            );

        border:
            1px solid rgba(65,255,120,0.12);

        border-radius: 18px;

        padding: 1.25rem;

        min-height: 135px;

        box-shadow:
            0 16px 40px
            rgba(0,0,0,0.18);
    }


    .card-label {
        font-size: 0.75rem;
        letter-spacing: 0.12em;
        text-transform: uppercase;
        color: #829389;
        font-weight: 700;
    }


    .card-value {
        font-size: 2.2rem;
        font-weight: 850;
        color: #ffffff;
        margin-top: 0.25rem;
    }


    .card-sub {
        font-size: 0.86rem;
        color: #a5b6aa;
        margin-top: 0.15rem;
    }


    .status-online {
        color: #35ef72;
        font-weight: 750;
    }


    .priority-high {
        color: #ff5e54;
        font-weight: 800;
    }


    .priority-medium {
        color: #f1c520;
        font-weight: 800;
    }


    .priority-low {
        color: #42e981;
        font-weight: 800;
    }


    .stButton > button {
        border-radius: 12px;

        border:
            1px solid rgba(50,255,110,0.22);

        background:
            linear-gradient(
                135deg,
                #10341e,
                #155e2c
            );

        color: white;

        font-weight: 700;

        transition:
            all 0.18s ease;

        box-shadow:
            0 8px 22px rgba(0,0,0,0.18);
    }


    .stButton > button:hover {
        transform: translateY(-1px);

        border-color: #38f873;

        box-shadow:
            0 0 24px
            rgba(52,255,112,0.15);
    }


    .stButton > button[kind="primary"] {
        background:
            linear-gradient(
                105deg,
                #18bd50 0%,
                #74ca28 48%,
                #f0c51e 100%
            );

        color: #061008;
        border: none;
        font-weight: 850;
    }


    div[data-testid="stChatMessage"] {
        background:
            rgba(13,28,18,0.82);

        border:
            1px solid rgba(255,255,255,0.05);

        border-radius: 18px;

        padding: 0.85rem 1rem;

        margin-bottom: 0.75rem;

        box-shadow:
            0 14px 30px
            rgba(0,0,0,0.14);
    }


    div[data-testid="stVerticalBlockBorderWrapper"] {
        background:
            rgba(12,27,18,0.74);

        border-radius: 18px;

        border:
            1px solid rgba(55,255,113,0.11);

        box-shadow:
            0 14px 30px
            rgba(0,0,0,0.14);
    }


    div[data-testid="metric-container"] {
        background:
            rgba(11,27,17,0.82);

        border:
            1px solid rgba(52,255,111,0.11);

        border-radius: 16px;

        padding: 1rem;

        box-shadow:
            0 10px 26px
            rgba(0,0,0,0.14);
    }


    div[data-baseweb="input"] > div,
    div[data-baseweb="textarea"] > div {
        background:
            rgba(11,25,16,0.92);

        border:
            1px solid rgba(50,255,110,0.14);

        border-radius: 14px;
    }


    input,
    textarea {
        color: white !important;
    }


    div[role="radiogroup"] label {
        background:
            rgba(12,28,18,0.65);

        border-radius: 10px;

        padding:
            0.48rem 0.7rem;

        margin-bottom:
            0.35rem;
    }


    div[data-testid="stAlert"] {
        border-radius: 14px;
    }


    hr {
        border-color:
            rgba(255,215,0,0.10);
    }


    ::-webkit-scrollbar {
        width: 9px;
    }


    ::-webkit-scrollbar-track {
        background: #070b08;
    }


    ::-webkit-scrollbar-thumb {
        background: #234f31;
        border-radius: 10px;
    }


    ::-webkit-scrollbar-thumb:hover {
        background: #317443;
    }

    </style>
    """,
    unsafe_allow_html=True,
)


# ==================================================
# SESSION STATE
# ==================================================

if "facts" not in st.session_state:
    st.session_state.facts = hent_json(
        FAKTAFIL,
        [],
    )


if "messages" not in st.session_state:
    st.session_state.messages = hent_json(
        CHATFIL,
        [],
    )


if "tasks" not in st.session_state:
    st.session_state.tasks = hent_json(
        OPGAVEFIL,
        [],
    )


if "memory_suggestion" not in st.session_state:
    st.session_state.memory_suggestion = None


if "task_suggestion" not in st.session_state:
    st.session_state.task_suggestion = None


if "task_notice" not in st.session_state:
    st.session_state.task_notice = None


if "workspace_page" not in st.session_state:
    st.session_state.workspace_page = "Dashboard"


if "last_page" not in st.session_state:
    st.session_state.last_page = None


# ==================================================
# HEADER
# ==================================================

st.markdown(
    '<div class="ak-brand">'
    'DOBELA · AK INTELLIGENCE'
    '</div>',
    unsafe_allow_html=True,
)


st.markdown(
    '<div class="ak-hero">'
    'AK-AI Partner'
    '</div>',
    unsafe_allow_html=True,
)


st.markdown(
    '<div class="ak-subtitle">'
    'AI-powered executive workspace'
    '</div>',
    unsafe_allow_html=True,
)


st.markdown(
    '<div class="reggae-line"></div>',
    unsafe_allow_html=True,
)


# ==================================================
# SIDEBAR
# ==================================================

with st.sidebar:

    st.markdown(
        "### 🌿 AK-AI"
    )

    st.caption(
        "Dobela Intelligence"
    )

    st.radio(
        "Workspace",
        [
            "Dashboard",
            "Chat",
            "Opgaver",
            "Mail Intelligence",
            "Hukommelse",
        ],
        key="workspace_page",
    )

    st.divider()

    aabne_count = len(
        [
            opgave
            for opgave in st.session_state.tasks
            if opgave["status"] == "ÅBEN"
        ]
    )

    hoeje_count = len(
        [
            opgave
            for opgave in st.session_state.tasks
            if (
                opgave["status"] == "ÅBEN"
                and opgave["prioritet"] == "HØJ"
            )
        ]
    )

    st.metric(
        "Åbne opgaver",
        aabne_count,
    )

    if hoeje_count:
        st.caption(
            f"🔴 {hoeje_count} med høj prioritet"
        )
    else:
        st.caption(
            "🟢 Ingen kritiske opgaver"
        )

    st.divider()

    st.caption(
        "● AI online"
    )


side = st.session_state.workspace_page


# ==================================================
# NAVIGATION / INTRO
# ==================================================

vis_hakuna_intro = (
    side == "Dashboard"
    and st.session_state.last_page != "Dashboard"
)

st.session_state.last_page = side


# ==================================================
# DASHBOARD
# ==================================================

if side == "Dashboard":

    if vis_hakuna_intro:

        components.html(
            """
            <style>

            html,
            body {
                margin: 0;
                padding: 0;
                overflow: hidden;
                background: transparent;
            }


            .party-stage {
                position: relative;
                width: 100%;
                height: 190px;
                overflow: hidden;
                display: flex;
                align-items: center;
                justify-content: center;
            }


            .party-text {
                position: absolute;
                white-space: nowrap;

                font-family:
                    Impact,
                    "Arial Black",
                    sans-serif;

                font-size:
                    clamp(
                        28px,
                        5vw,
                        68px
                    );

                font-weight: 900;
                font-style: italic;
                letter-spacing: 2px;

                background:
                    linear-gradient(
                        90deg,
                        #20e65d,
                        #ffe13b,
                        #ff493d,
                        #ff52dc,
                        #45d9ff,
                        #20e65d
                    );

                background-size:
                    300% 100%;

                -webkit-background-clip:
                    text;

                background-clip:
                    text;

                color:
                    transparent;

                filter:
                    drop-shadow(
                        0 0 8px
                        rgba(
                            255,
                            255,
                            255,
                            0.30
                        )
                    )
                    drop-shadow(
                        0 0 18px
                        rgba(
                            40,
                            255,
                            100,
                            0.35
                        )
                    );

                animation:
                    flyAcross
                    3.8s
                    cubic-bezier(
                        .15,
                        .8,
                        .25,
                        1
                    )
                    forwards,

                    rainbow
                    1.1s
                    linear
                    infinite,

                    partyShake
                    0.20s
                    ease-in-out
                    infinite
                    alternate;
            }


            .glow {
                position: absolute;

                width: 430px;
                height: 110px;

                border-radius: 50%;

                background:
                    radial-gradient(
                        ellipse,
                        rgba(
                            44,
                            255,
                            104,
                            0.25
                        ),
                        rgba(
                            255,
                            211,
                            35,
                            0.12
                        ),
                        transparent
                        70%
                    );

                filter: blur(22px);

                animation:
                    glowPulse
                    0.6s
                    ease-in-out
                    infinite
                    alternate;
            }


            .spark {
                position: absolute;

                width: 9px;
                height: 9px;

                border-radius: 50%;

                background: #ffe338;

                box-shadow:
                    0 0 14px
                    #ffe338;

                opacity: 0;

                animation:
                    sparkle
                    2.5s
                    ease-out
                    forwards;
            }


            .s1 {
                left: 12%;
                top: 30%;
                animation-delay: .25s;
            }


            .s2 {
                left: 27%;
                top: 72%;
                animation-delay: .6s;
                background: #35ff75;
            }


            .s3 {
                left: 48%;
                top: 18%;
                animation-delay: .9s;
                background: #ff4c42;
            }


            .s4 {
                left: 69%;
                top: 70%;
                animation-delay: 1.25s;
                background: #45d9ff;
            }


            .s5 {
                left: 87%;
                top: 31%;
                animation-delay: 1.6s;
                background: #ff52dc;
            }


            .s6 {
                left: 38%;
                top: 52%;
                animation-delay: 1.9s;
                background: #ffe338;
            }


            @keyframes flyAcross {

                0% {
                    transform:
                        translateX(
                            -150vw
                        )
                        scale(.40)
                        rotate(-9deg);

                    opacity: 0;
                }


                18% {
                    opacity: 1;
                }


                45% {
                    transform:
                        translateX(0)
                        scale(1.15)
                        rotate(1deg);

                    opacity: 1;
                }


                65% {
                    transform:
                        translateX(0)
                        scale(1)
                        rotate(-1deg);

                    opacity: 1;
                }


                100% {
                    transform:
                        translateX(
                            150vw
                        )
                        scale(.50)
                        rotate(9deg);

                    opacity: 0;
                }
            }


            @keyframes rainbow {

                0% {
                    background-position:
                        0% 50%;
                }


                100% {
                    background-position:
                        300% 50%;
                }
            }


            @keyframes partyShake {

                from {
                    margin-top: -3px;
                }


                to {
                    margin-top: 3px;
                }
            }


            @keyframes glowPulse {

                from {
                    transform:
                        scale(.8);

                    opacity: .35;
                }


                to {
                    transform:
                        scale(1.25);

                    opacity: .9;
                }
            }


            @keyframes sparkle {

                0% {
                    opacity: 0;

                    transform:
                        scale(.2)
                        translateY(
                            35px
                        );
                }


                25% {
                    opacity: 1;

                    transform:
                        scale(1.8)
                        translateY(0);
                }


                100% {
                    opacity: 0;

                    transform:
                        scale(.1)
                        translateY(
                            -85px
                        );
                }
            }

            </style>


            <div class="party-stage">

                <div class="glow"></div>

                <div class="spark s1"></div>
                <div class="spark s2"></div>
                <div class="spark s3"></div>
                <div class="spark s4"></div>
                <div class="spark s5"></div>
                <div class="spark s6"></div>

                <div class="party-text">
                    HAKUNA MATATA — LET'S FUCKING GO!
                </div>

            </div>
            """,
            height=190,
        )


    st.header(
        "Executive Cockpit"
    )


    st.caption(
        "Det vigtigste lige nu"
    )


    aabne = [
        opgave
        for opgave in st.session_state.tasks
        if opgave["status"] == "ÅBEN"
    ]


    afsluttede = [
        opgave
        for opgave in st.session_state.tasks
        if opgave["status"] == "AFSLUTTET"
    ]


    hoeje = [
        opgave
        for opgave in aabne
        if opgave["prioritet"] == "HØJ"
    ]


    col1, col2, col3, col4 = st.columns(4)


    with col1:

        kort_html = (
            '<div class="cockpit-card">'
            '<div class="card-label">Åbne opgaver</div>'
            f'<div class="card-value">{len(aabne)}</div>'
            '<div class="card-sub">Aktive handlinger</div>'
            '</div>'
        )

        st.markdown(
            kort_html,
            unsafe_allow_html=True,
        )


    with col2:

        kort_html = (
            '<div class="cockpit-card">'
            '<div class="card-label">Høj prioritet</div>'
            f'<div class="card-value">{len(hoeje)}</div>'
            '<div class="card-sub">Kræver opmærksomhed</div>'
            '</div>'
        )

        st.markdown(
            kort_html,
            unsafe_allow_html=True,
        )


    with col3:

        kort_html = (
            '<div class="cockpit-card">'
            '<div class="card-label">Afsluttede</div>'
            f'<div class="card-value">{len(afsluttede)}</div>'
            '<div class="card-sub">Gennemførte opgaver</div>'
            '</div>'
        )

        st.markdown(
            kort_html,
            unsafe_allow_html=True,
        )


    with col4:

        kort_html = (
            '<div class="cockpit-card">'
            '<div class="card-label">Hukommelse</div>'
            f'<div class="card-value">{len(st.session_state.facts)}</div>'
            '<div class="card-sub">Kendte fakta</div>'
            '</div>'
        )

        st.markdown(
            kort_html,
            unsafe_allow_html=True,
        )


    st.write("")
    st.write("")


    col1, col2 = st.columns(
        [2, 1]
    )


    with col1:

        st.subheader(
            "Prioriterede opgaver"
        )


        prioriterede = sorted(
            aabne,
            key=prioritet_sortering,
        )[:5]


        if not prioriterede:

            st.success(
                "Ingen åbne opgaver lige nu."
            )


        for opgave in prioriterede:

            with st.container(
                border=True
            ):

                top1, top2 = st.columns(
                    [5, 1]
                )


                with top1:

                    st.markdown(
                        f"### {opgave['titel']}"
                    )


                with top2:

                    prioritet = opgave[
                        "prioritet"
                    ]


                    if prioritet == "HØJ":

                        st.markdown(
                            '<span class="priority-high">'
                            '● HØJ'
                            '</span>',
                            unsafe_allow_html=True,
                        )


                    elif prioritet == "MELLEM":

                        st.markdown(
                            '<span class="priority-medium">'
                            '● MELLEM'
                            '</span>',
                            unsafe_allow_html=True,
                        )


                    else:

                        st.markdown(
                            '<span class="priority-low">'
                            '● LAV'
                            '</span>',
                            unsafe_allow_html=True,
                        )


                st.caption(
                    f"Ansvarlig: "
                    f"{opgave['ansvarlig']} · "
                    f"Deadline: "
                    f"{opgave['deadline']}"
                )


                if opgave["beskrivelse"]:

                    st.write(
                        opgave["beskrivelse"]
                    )


    with col2:

        st.subheader(
            "Quick access"
        )


        if st.button(
            "💬 Åbn AI-chat",
            use_container_width=True,
        ):

            st.session_state.workspace_page = "Chat"

            st.rerun()


        if st.button(
            "✓ Se alle opgaver",
            use_container_width=True,
        ):

            st.session_state.workspace_page = "Opgaver"

            st.rerun()


        if st.button(
            "✉ Mail Intelligence",
            use_container_width=True,
        ):

            st.session_state.workspace_page = (
                "Mail Intelligence"
            )

            st.rerun()


        if st.button(
            "🧠 Hukommelse",
            use_container_width=True,
        ):

            st.session_state.workspace_page = "Hukommelse"

            st.rerun()


        st.write("")


        with st.container(
            border=True
        ):

            st.markdown(
                "### System"
            )


            st.markdown(
                '<span class="status-online">'
                '● AI online'
                '</span>',
                unsafe_allow_html=True,
            )


            st.caption(
                "Chat · Tasks · Documents · Memory"
            )


            st.caption(
                "Outlook integration: næste fase"
            )


# ==================================================
# CHAT
# ==================================================

elif side == "Chat":

    st.header(
        "AI Workspace"
    )


    st.caption(
        "Tal naturligt med AK-AI-Partner"
    )


    if st.session_state.task_notice:

        st.success(
            st.session_state.task_notice
        )

        st.session_state.task_notice = None


    if st.session_state.task_suggestion:

        forslag = (
            st.session_state.task_suggestion
        )


        st.warning(
            "Det her ligner en mulig opgave. "
            "Skal jeg oprette den?"
        )


        with st.container(
            border=True
        ):

            st.markdown(
                f"### {forslag['titel']}"
            )


            st.write(
                f"**Ansvarlig:** "
                f"{forslag['ansvarlig']}"
            )


            st.write(
                f"**Deadline:** "
                f"{forslag['deadline']}"
            )


            st.write(
                f"**Prioritet:** "
                f"{forslag['prioritet']}"
            )


            col1, col2 = st.columns(2)


            with col1:

                if st.button(
                    "Opret opgave",
                    type="primary",
                    use_container_width=True,
                ):

                    opgave = opret_ai_opgave(
                        forslag,
                        "Chat - godkendt",
                    )


                    st.session_state.task_suggestion = None


                    if opgave:

                        st.session_state.task_notice = (
                            f"Opgaven "
                            f"'{opgave['titel']}' "
                            f"er oprettet."
                        )


                    else:

                        st.session_state.task_notice = (
                            "Opgaven findes allerede."
                        )


                    st.rerun()


            with col2:

                if st.button(
                    "Nej tak",
                    key="nej_opgave",
                    use_container_width=True,
                ):

                    st.session_state.task_suggestion = None

                    st.rerun()


    if st.session_state.memory_suggestion:

        st.info(
            "Jeg foreslår at huske:\n\n"
            + st.session_state.memory_suggestion
        )


        col1, col2 = st.columns(2)


        with col1:

            if st.button(
                "Husk det",
                key="husk_det",
            ):

                forslag = (
                    st.session_state.memory_suggestion
                )


                if (
                    forslag
                    not in st.session_state.facts
                ):

                    st.session_state.facts.append(
                        forslag
                    )


                    gem_json(
                        FAKTAFIL,
                        st.session_state.facts,
                    )


                st.session_state.memory_suggestion = None

                st.rerun()


        with col2:

            if st.button(
                "Nej tak",
                key="nej_hukommelse",
            ):

                st.session_state.memory_suggestion = None

                st.rerun()


    aabne_opgaver = [
        opgave
        for opgave in st.session_state.tasks
        if opgave["status"] == "ÅBEN"
    ]


    opgavekontekst = "\n".join(
        (
            f"- {opgave['titel']} | "
            f"Ansvarlig: {opgave['ansvarlig']} | "
            f"Deadline: {opgave['deadline']} | "
            f"Prioritet: {opgave['prioritet']}"
        )
        for opgave in aabne_opgaver
    )


    for message in st.session_state.messages:

        with st.chat_message(
            message["role"]
        ):

            st.write(
                message["content"]
            )


    besked = st.chat_input(
        "Skriv til AK-AI-Partner..."
    )


    if besked:

        st.session_state.task_suggestion = None


        st.session_state.messages.append(
            {
                "role": "user",
                "content": besked,
            }
        )


        gem_json(
            CHATFIL,
            st.session_state.messages,
        )


        kendte_fakta = "\n".join(
            f"- {faktum}"
            for faktum
            in st.session_state.facts
        )


        response = client.responses.create(
            model="gpt-5.4-mini",

            instructions=(
                f"{partnerprofil}\n\n"

                "Langtidshukommelse:\n"
                f"{kendte_fakta}\n\n"

                "Aktuelle åbne opgaver:\n"
                f"{opgavekontekst}\n\n"

                "Brug oplysningerne som kontekst. "
                "Vær konkret, proaktiv og praktisk."
            ),

            input=st.session_state.messages,
        )


        svar = response.output_text


        st.session_state.messages.append(
            {
                "role": "assistant",
                "content": svar,
            }
        )


        gem_json(
            CHATFIL,
            st.session_state.messages,
        )


        opgaveanalyse = analyser_for_opgave(
            besked
        )


        if (
            opgaveanalyse["status"] == "CLEAR"
            and opgaveanalyse["titel"]
        ):

            st.session_state.task_suggestion = None


            opgave = opret_ai_opgave(
                opgaveanalyse,
                "Chat - automatisk",
            )


            if opgave:

                st.session_state.task_notice = (
                    f"Jeg oprettede automatisk "
                    f"opgaven "
                    f"'{opgave['titel']}'."
                )


        elif (
            opgaveanalyse["status"]
            == "UNCERTAIN"
            and opgaveanalyse["titel"]
        ):

            st.session_state.task_suggestion = (
                opgaveanalyse
            )


        else:

            st.session_state.task_suggestion = None


        forslag = foreslaa_hukommelse(
            besked
        )


        if (
            forslag
            and forslag
            not in st.session_state.facts
        ):

            st.session_state.memory_suggestion = forslag


        st.rerun()


# ==================================================
# OPGAVER
# ==================================================

elif side == "Opgaver":

    st.header(
        "Task Intelligence"
    )


    st.caption(
        "Opgaver, ansvar og deadlines"
    )


    aabne = [
        opgave
        for opgave in st.session_state.tasks
        if opgave["status"] == "ÅBEN"
    ]


    afsluttede = [
        opgave
        for opgave in st.session_state.tasks
        if opgave["status"] == "AFSLUTTET"
    ]


    hoeje = [
        opgave
        for opgave in aabne
        if opgave["prioritet"] == "HØJ"
    ]


    col1, col2, col3 = st.columns(3)


    col1.metric(
        "Åbne",
        len(aabne),
    )


    col2.metric(
        "Høj prioritet",
        len(hoeje),
    )


    col3.metric(
        "Afsluttede",
        len(afsluttede),
    )


    with st.expander(
        "＋ Opret ny opgave"
    ):

        with st.form(
            "ny_opgave"
        ):

            titel = st.text_input(
                "Opgave"
            )


            beskrivelse = st.text_area(
                "Beskrivelse"
            )


            col1, col2 = st.columns(2)


            with col1:

                ansvarlig = st.text_input(
                    "Ansvarlig",
                    placeholder=(
                        "Fx KE eller Antoaneta"
                    ),
                )


            with col2:

                deadline = st.text_input(
                    "Deadline",
                    placeholder=(
                        "Fx 20-08-2026"
                    ),
                )


            prioritet = st.selectbox(
                "Prioritet",
                [
                    "HØJ",
                    "MELLEM",
                    "LAV",
                ],
                index=1,
            )


            opret = st.form_submit_button(
                "Opret opgave",
                type="primary",
            )


            if opret:

                if titel.strip():

                    ny = ny_opgave(
                        titel=titel,
                        beskrivelse=beskrivelse,
                        ansvarlig=ansvarlig,
                        deadline=deadline,
                        prioritet=prioritet,
                    )


                    if not opgave_findes(
                        ny["titel"],
                        ny["ansvarlig"],
                        ny["deadline"],
                    ):

                        st.session_state.tasks.append(
                            ny
                        )


                        gem_opgaver()


                        st.rerun()


                    else:

                        st.warning(
                            "En tilsvarende "
                            "åben opgave findes allerede."
                        )


                else:

                    st.warning(
                        "Opgaven skal have en titel."
                    )


    st.write("")


    aabne = sorted(
        aabne,
        key=prioritet_sortering,
    )


    if not aabne:

        st.success(
            "Ingen åbne opgaver."
        )


    for opgave in aabne:

        with st.container(
            border=True
        ):

            top1, top2 = st.columns(
                [5, 1]
            )


            with top1:

                st.subheader(
                    opgave["titel"]
                )


            with top2:

                if opgave["prioritet"] == "HØJ":

                    st.markdown(
                        '<span class="priority-high">'
                        '● HØJ'
                        '</span>',
                        unsafe_allow_html=True,
                    )


                elif opgave["prioritet"] == "MELLEM":

                    st.markdown(
                        '<span class="priority-medium">'
                        '● MELLEM'
                        '</span>',
                        unsafe_allow_html=True,
                    )


                else:

                    st.markdown(
                        '<span class="priority-low">'
                        '● LAV'
                        '</span>',
                        unsafe_allow_html=True,
                    )


            if opgave["beskrivelse"]:

                st.write(
                    opgave["beskrivelse"]
                )


            col1, col2 = st.columns(2)


            with col1:

                st.write(
                    f"**Ansvarlig**  \n"
                    f"{opgave['ansvarlig']}"
                )


            with col2:

                st.write(
                    f"**Deadline**  \n"
                    f"{opgave['deadline']}"
                )


            st.caption(
                f"Kilde: "
                f"{opgave['kilde']}"
            )


            col1, col2 = st.columns(2)


            with col1:

                if st.button(
                    "✓ Afslut",
                    key=f"afslut_{opgave['id']}",
                    use_container_width=True,
                ):

                    afslut_opgave(
                        opgave["id"]
                    )


                    st.rerun()


            with col2:

                if st.button(
                    "Slet",
                    key=f"slet_task_{opgave['id']}",
                    use_container_width=True,
                ):

                    slet_opgave(
                        opgave["id"]
                    )


                    st.rerun()


    with st.expander(
        f"Afsluttede opgaver "
        f"({len(afsluttede)})"
    ):

        for opgave in afsluttede:

            st.write(
                f"**{opgave['titel']}**"
            )


            st.caption(
                f"{opgave['ansvarlig']} · "
                f"{opgave['deadline']}"
            )


            col1, col2 = st.columns(2)


            with col1:

                if st.button(
                    "Genåbn",
                    key=f"genaabn_{opgave['id']}",
                ):

                    genaabn_opgave(
                        opgave["id"]
                    )


                    st.rerun()


            with col2:

                if st.button(
                    "Slet permanent",
                    key=f"slet_done_{opgave['id']}",
                ):

                    slet_opgave(
                        opgave["id"]
                    )


                    st.rerun()


# ==================================================
# MAIL INTELLIGENCE
# ==================================================

elif side == "Mail Intelligence":

    st.header(
        "Mail Intelligence"
    )


    st.caption(
        "Email + dokumentanalyse i én arbejdsgang"
    )


    mailtekst = st.text_area(
        "Email",
        height=280,
        placeholder=(
            "Indsæt emailens tekst her..."
        ),
    )


    uploaded_files = st.file_uploader(
        "Vedhæftninger",
        type=[
            "pdf",
            "docx",
            "xlsx",
            "pptx",
            "txt",
        ],
        accept_multiple_files=True,
    )


    dokumenter = []


    if uploaded_files:

        for fil in uploaded_files:

            try:

                tekst = laes_vedhaeftning(
                    fil
                )


                dokumenter.append(
                    {
                        "navn": fil.name,
                        "tekst": tekst,
                    }
                )


                st.success(
                    f"✓ {fil.name} er klar"
                )


            except Exception as fejl:

                st.error(
                    f"Kunne ikke læse "
                    f"{fil.name}: {fejl}"
                )


    if st.button(
        "Analyser materiale",
        type="primary",
    ):

        if (
            not mailtekst.strip()
            and not dokumenter
        ):

            st.warning(
                "Indsæt en mail "
                "eller mindst én vedhæftning."
            )


        else:

            with st.spinner(
                "AK-AI analyserer materialet..."
            ):

                analyse = analyser_mail(
                    mailtekst,
                    dokumenter,
                )


            st.subheader(
                "Executive analyse"
            )


            st.markdown(
                analyse
            )


# ==================================================
# HUKOMMELSE
# ==================================================

elif side == "Hukommelse":

    st.header(
        "Memory Intelligence"
    )


    st.caption(
        "Det AK-AI aktivt husker og bruger som kontekst"
    )


    ny_fakta = st.text_input(
        "Tilføj hukommelse"
    )


    if st.button(
        "Gem",
        type="primary",
    ):

        if ny_fakta.strip():

            if (
                ny_fakta.strip()
                not in st.session_state.facts
            ):

                st.session_state.facts.append(
                    ny_fakta.strip()
                )


                gem_json(
                    FAKTAFIL,
                    st.session_state.facts,
                )


            st.rerun()


    if not st.session_state.facts:

        st.info(
            "Ingen gemte oplysninger endnu."
        )


    for i, faktum in enumerate(
        st.session_state.facts
    ):

        with st.container(
            border=True
        ):

            col1, col2 = st.columns(
                [6, 1]
            )


            with col1:

                st.write(
                    faktum
                )


            with col2:

                if st.button(
                    "Slet",
                    key=f"slet_fact_{i}",
                ):

                    st.session_state.facts.pop(
                        i
                    )


                    gem_json(
                        FAKTAFIL,
                        st.session_state.facts,
                    )


                    st.rerun()